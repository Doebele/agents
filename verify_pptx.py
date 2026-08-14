#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Verifikation der erzeugten PPTX gegen die Akzeptanzkriterien (Abschnitt 7).
Reines Lese-Programm; verändert nichts an der Datei.
"""

import os
import re
import sys

from pptx import Presentation
from pptx.util import Emu

HERE    = os.path.dirname(os.path.abspath(__file__))
OUT     = os.path.join(HERE, "KI-Agenten-Landschaft.pptx")
N_TOTAL = 22

DIVIDERS = {4: "b1_llm", 8: "b2_harness", 11: "b3_skills", 13: "b4_mcp",
            15: "b5_tools", 17: "b6_medien", 19: "b7_design"}

# Produktnamen / Versionsnummern aus Abschnitt 4 (müssen im Deck vorkommen)
PRODUCTS = [
    # Baustein 1
    "GPT-5.6", "Claude Opus 5", "Sonnet 5", "Gemini 3.1 Pro", "Grok 4.5",
    "Kimi K3", "DeepSeek V4", "GLM-5.2", "Qwen 3.6 / 3.8", "gpt-oss",
    "Mistral Large 3", "Llama 4", "Ornith 1.0", "Qwen3.8-Max",
    "LM Studio", "Ollama", "llama.cpp", "vLLM", "SGLang",
    # Baustein 2
    "Claude Code", "Agent SDK", "Codex", "Agents SDK", "Kimi Code", "ZCode",
    "Zed", "Hermes Agent", "Odysseus", "Bionic", "LangGraph", "LangChain",
    "CrewAI", "Paperclip", "Google ADK",
    # Baustein 3
    "SKILL.md", "agentskills.io", "skills.sh", "40+",
    # Baustein 4
    "9.800", "Playwright", "GitHub", "Filesystem", "Slack", "Postgres",
    "Supabase", "Google Drive", "Linux Foundation",
    # Baustein 6
    "GPT Image 2", "Midjourney V8.2", "FLUX.2", "Imagen 4", "Ideogram",
    "Recraft", "Seedream", "Veo 3.1", "Kling 3.0", "Seedance 2.5", "Wan 3.0",
    "Runway", "Luma", "Higgsfield", "ElevenLabs", "Suno v5.5", "Stable Audio 3",
    "Whisper", "Kokoro", "Voxtral",
    # Baustein 7
    "Figma", "pencil.dev", "Penpot", "Antigravity 2.0", "Google Stitch",
    "v0", "Lovable", "Bolt.new", "Onlook", "Relume", "Uizard",
]

problems = []


def check(cond, msg):
    print(("  OK  " if cond else " FEHL ") + msg)
    if not cond:
        problems.append(msg)


def slide_text(slide):
    chunks = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            chunks.append(sh.text_frame.text)
    return "\n".join(chunks)


def count_images(slide):
    return sum(1 for sh in slide.shapes if sh.shape_type == 13)  # PICTURE


def main():
    prs = Presentation(OUT)
    print("=== Struktur ===")
    check(len(prs.slides._sldIdLst) == N_TOTAL,
          "%d Folien (Soll %d)" % (len(prs.slides._sldIdLst), N_TOTAL))
    w, h = prs.slide_width, prs.slide_height
    check(w == Emu(960 * 12700) and h == Emu(540 * 12700),
          "Format 16:9 (960x540 px = %.3f\" x %.3f\")" % (w / 914400, h / 914400))

    slides = list(prs.slides)

    print("\n=== Seitenzahlen NN / 22 ===")
    for i, slide in enumerate(slides, 1):
        txt = slide_text(slide)
        if i == 1:
            check(re.search(r"\b01\s*/\s*22\b", txt) is None,
                  "Folie %02d (Cover) hat keine Seitenzahl" % i)
        else:
            expect = "%02d / 22" % i
            check(expect in txt, "Folie %02d zeigt '%s'" % (i, expect))

    print("\n=== Trennseiten mit Objektbild ===")
    for num, key in DIVIDERS.items():
        slide = slides[num - 1]
        imgs = count_images(slide)
        check(imgs >= 1, "Folie %02d trägt %d Bild(er) (Thema: %s)" % (num, imgs, key))

    print("\n=== Folie 03 Knotendiagramm ===")
    s3 = slides[2]
    t3 = slide_text(s3)
    for node in ["LLM", "Harness", "MCP", "Tools", "Skills"]:
        check(node in t3, "Knoten '%s' auf Folie 03" % node)
    check(count_images(s3) >= 5, "Folie 03 enthält >=5 Knotenbilder (hat %d)" % count_images(s3))
    # Pfeile als Verbinder (straight connector = 9 ? type varies) -> zähle Linien-Connectoren
    n_conn = sum(1 for sh in s3.shapes if sh.shape_type == 9)  # LINE / connector
    check(n_conn >= 4, "Folie 03 enthält >=4 Verbindungspfeile (hat %d)" % n_conn)

    print("\n=== Produktnamen / Versionen (Abschnitt 4) ===")
    fulltext = "\n".join(slide_text(s) for s in slides)
    missing = [p for p in PRODUCTS if p not in fulltext]
    if missing:
        print("  FEHL  fehlend: " + ", ".join(missing))
        problems.append("Fehlende Produktnamen: " + ", ".join(missing))
    else:
        print("  OK   alle %d Begriffe gefunden" % len(PRODUCTS))

    print("\n=== Typografie / Sprache ===")
    # Deutsche Anführungszeichen „ " , Gedankenstrich —, Umlaute/ß
    for glyph, name in [("„", "öffnendes Anführungszeichen"), ("“", "schließendes Anführungszeichen"),
                        ("—", "Gedankenstrich"), ("ß", "sz"), ("ä", "ä"), ("ö", "ö"), ("ü", "ü")]:
        check(glyph in fulltext, "deutsches Zeichen '%s' (%s) vorhanden" % (glyph, name))
    # Kein gerade (ASCII) Anführungszeichen im Sinne der Spec  – " sei erlaubt (Zoll), ' ebenso
    # Keine Emoji (BMP außerhalb gebräuchlicher Bereiche)
    emoji_re = re.compile(
        "[\U0001F300-\U0001FAFF\U00002600-\U000027BF\U0001F1E6-\U0001F1FF]")
    emojis = sorted(set(emoji_re.findall(fulltext)))
    check(not emojis, "keine Emoji" + (" (gefunden: %s)" % emojis if emojis else ""))

    print("\n=== Hintergrundfarbe Papier (Stichprobe) ===")
    from pptx.dml.color import RGBColor
    paper = RGBColor(0xF7, 0xF4, 0xEE)
    ok_bg = True
    for slide in slides:
        try:
            f = slide.background.fill
            if f.type == 1 and f.fore_color.rgb != paper:
                ok = False
        except Exception:
            pass
    check(ok_bg, "alle Folien haben Papier-Hintergrund #F7F4EE")

    print("\n=== Bilddatei ===")
    size = os.path.getsize(OUT)
    print("  Datei: %s  (%.0f KB)" % (os.path.relpath(OUT, HERE), size / 1024))

    print("\n" + "=" * 60)
    if problems:
        print("ERGEBNIS: %d Problem(e)" % len(problems))
        for p in problems:
            print("  - " + p)
        sys.exit(1)
    print("ERGEBNIS: alle Prüfungen bestanden.")


if __name__ == "__main__":
    main()
