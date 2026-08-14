#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Renderer: Referenz-Präsentation (.page/.pptd, Dieter-Rams-Designsystem) -> PPTX (16:9).

Liest die schreibgeschützten Referenzdateien unter ../ki-landschaft-praesentation/
und erzeugt eine neue PPTX-Datei unter output/. Die Quelldateien bleiben unangetastet.

Maßstab:  Referenz-Canvas 960 x 540 px  ==  PPTX-Folie 13,333" x 7,5" (16:9).
         1 px  ==  12700 EMU  (960*12700 = 12.192.000 ; 540*12700 = 6.858.000).
"""

import os
import re
import sys
import yaml

from PIL import Image as PILImage

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------- #
#  Pfade
# --------------------------------------------------------------------------- #
HERE      = os.path.dirname(os.path.abspath(__file__))
REF_DIR   = os.path.abspath(os.path.join(HERE, "..", "ki-landschaft-praesentation"))
PPTD_PATH = os.path.join(REF_DIR, "ki-landschaft.pptd")
OUT_PPTX  = os.path.join(HERE, "KI-Agenten-Landschaft.pptx")

# --------------------------------------------------------------------------- #
#  Maßstab-Helper
# --------------------------------------------------------------------------- #
def E(v):
    """px-Wert (960x540-Canvas) -> EMU."""
    return Emu(int(round(v * 12700)))

# --------------------------------------------------------------------------- #
#  Designsystem (aus .pptd theme)
# --------------------------------------------------------------------------- #
COL_TOKENS = {
    "$paper":    RGBColor(0xF7, 0xF4, 0xEE),
    "$ink":      RGBColor(0x1C, 0x1B, 0x18),
    "$muted":    RGBColor(0x6E, 0x6A, 0x60),
    "$accent":   RGBColor(0xD9, 0x50, 0x0B),
    "$hairline": RGBColor(0xDA, 0xD5, 0xC8),
    "$faint":    RGBColor(0xE9, 0xE4, 0xD6),
}

# Textstil-Vorgaben (theme.textStyles)
STYLES = {
    "$kicker": dict(font="JetBrains Mono", size=11, color=COL_TOKENS["$accent"],
                    spc=200, bold=False, line_height=None),
    "$title":  dict(font="Liter", size=30, color=COL_TOKENS["$ink"],
                    bold=True, line_height=None),
    "$body":   dict(font="Liter", size=15, color=COL_TOKENS["$ink"],
                    bold=False, line_height=1.65),
    "$gloss":  dict(font="Liter", size=13, color=COL_TOKENS["$muted"],
                    bold=False, line_height=1.5),
    "$label":  dict(font="JetBrains Mono", size=10, color=COL_TOKENS["$muted"],
                    spc=150, bold=False, line_height=None),
}

ALIGN_H = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}
ALIGN_V = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


# --------------------------------------------------------------------------- #
#  Rich-Text-Parser  (Mini-HTML: <p>, <strong>, <span style="color:$…">)
# --------------------------------------------------------------------------- #
_TAG = re.compile(r"(<strong>.*?</strong>|<span[^>]*>.*?</span>)", re.S)


def _parse_runs(s):
    """Ein Paragraph-String -> Liste (text, attrs)."""
    runs = []
    pos = 0
    for m in _TAG.finditer(s):
        if m.start() > pos:
            runs.append((s[pos:m.start()], {}))
        tag = m.group(0)
        if tag.startswith("<strong"):
            inner = re.sub(r"</?strong>", "", tag)
            runs.append((inner, {"bold": True}))
        else:  # <span ...>
            inner = re.search(r">(.*?)</span>", tag, re.S).group(1)
            color = None
            if "color:$accent" in tag:
                color = COL_TOKENS["$accent"]
            elif "color:$muted" in tag:
                color = COL_TOKENS["$muted"]
            attrs = {}
            if color is not None:
                attrs["color"] = color
            runs.append((inner, attrs))
        pos = m.end()
    if pos < len(s):
        runs.append((s[pos:], {}))
    runs = [(t, a) for (t, a) in runs if t != ""]
    return runs or [("", {})]


def _parse_text(text):
    """content.text -> Liste von Paragraphen (je Liste von Runs)."""
    if text is None:
        return [[]]
    text = str(text).strip("\n")
    if "<p" in text:
        paras = re.findall(r"<p[^>]*>(.*?)</p>", text, flags=re.S)
        return [_parse_runs(p) for p in paras]
    # mehrzeiliger Klartext -> mehrere Paragraphen
    lines = text.split("\n")
    if len(lines) > 1:
        return [_parse_runs(l) for l in lines]
    return [_parse_runs(text)]


# --------------------------------------------------------------------------- #
#  Element-Primitive
# --------------------------------------------------------------------------- #
def _resolve_color(c):
    if c is None:
        return None
    if isinstance(c, str):
        return COL_TOKENS.get(c)
    return c


def _zero_margins(tf):
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0


def add_text(slide, bounds, content):
    """content ist der content-Block eines Text-Elements."""
    x, y, w, h = (E(v) for v in bounds)

    # Basis-Stil aus style-Referenz
    style_ref = content.get("style")
    base = dict(STYLES.get(style_ref, {}))  # Kopie

    # Einzelfeld-Overrides
    if "fontSize" in content:
        base["size"] = content["fontSize"]
    if "fontFamily" in content:
        base["font"] = content["fontFamily"]
    if "color" in content:
        base["color"] = _resolve_color(content["color"])
    if "bold" in content:
        base["bold"] = content["bold"]
    if "lineHeight" in content:
        base["line_height"] = content["lineHeight"]
    if "letterSpacing" in content:
        base["spc"] = int(content["letterSpacing"] * 100)

    align = content.get("align", ["left", "top"])
    align_h = align[0] if isinstance(align, (list, tuple)) else align
    align_v = align[1] if isinstance(align, (list, tuple)) and len(align) > 1 else "top"

    wrap = content.get("wrap", True)

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = ALIGN_V.get(align_v, MSO_ANCHOR.TOP)
    _zero_margins(tf)

    paragraphs = _parse_text(content.get("text", ""))
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN_H.get(align_h, PP_ALIGN.LEFT)
        if base.get("line_height"):
            p.line_spacing = base["line_height"]
        for (rt, ra) in runs:
            r = p.add_run()
            r.text = rt
            f = r.font
            f.size = Pt(base.get("size", 15))
            f.name = base.get("font", "Liter")
            f.bold = ra.get("bold", base.get("bold", False))
            f.color.rgb = ra.get("color", base.get("color", COL_TOKENS["$ink"]))
            # Laufoerweite (Buchstaben-Abstand)
            rPr = r._r.get_or_add_rPr()
            if base.get("spc"):
                rPr.set("spc", str(int(base["spc"])))
    return tb


def add_rect(slide, bounds, fill):
    x, y, w, h = (E(v) for v in bounds)
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.line.fill.background()  # keine Kontur
    color = _resolve_color(fill.get("color"))
    if color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    return shp


def _arrow_end(ln_elem, kind, present):
    """headEnd/tailEnd sauber setzen (Schema-Reihenfolge)."""
    tag = "a:headEnd" if kind == "head" else "a:tailEnd"
    for e in ln_elem.findall(qn(tag)):
        ln_elem.remove(e)
    if present:
        el = ln_elem.makeelement(qn(tag), {"type": "triangle", "w": "med", "len": "med"})
        ln_elem.append(el)


def add_line(slide, el):
    bounds = el["bounds"]
    x, y, w, h = bounds
    vb = el.get("viewBox", [w, h])
    vw, vh = vb
    pts_raw = el["points"]
    if isinstance(pts_raw, str):
        coords = [tuple(map(float, p.split(","))) for p in pts_raw.split()]
    else:
        coords = [tuple(map(float, p)) for p in pts_raw]
    # viewBox -> absolute Slide-Koordinaten
    abs_pts = []
    for (px, py) in coords:
        ax = x + (px * w / vw if vw else 0)
        ay = y + (py * h / vh if vh else 0)
        abs_pts.append((ax, ay))

    border = el.get("border", {})
    width_pt = border.get("width", 1.5)
    color = _resolve_color(border.get("color", COL_TOKENS["$muted"]))
    arrow = el.get("arrow", [None, None])
    head = arrow[0] == "arrow" if isinstance(arrow, (list, tuple)) else False
    tail = (arrow[1] == "arrow") if isinstance(arrow, (list, tuple)) and len(arrow) > 1 else False

    if len(abs_pts) == 2:
        (x1, y1), (x2, y2) = abs_pts
        cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                        E(x1), E(y1), E(x2), E(y2))
        cn.line.color.rgb = color
        cn.line.width = Pt(width_pt)
        ln = cn.line._get_or_add_ln()
        _arrow_end(ln, "head", head)
        _arrow_end(ln, "tail", tail)
    else:
        # Polyline aus Segmenten, Pfeilkopf nur am letzten Segment
        for i in range(len(abs_pts) - 1):
            (x1, y1), (x2, y2) = abs_pts[i], abs_pts[i + 1]
            cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                            E(x1), E(y1), E(x2), E(y2))
            cn.line.color.rgb = color
            cn.line.width = Pt(width_pt)
            ln = cn.line._get_or_add_ln()
            _arrow_end(ln, "head", False)
            _arrow_end(ln, "tail", tail and i == len(abs_pts) - 2)


def add_image(slide, el, base_dir):
    bounds = el["bounds"]
    x, y, w, h = bounds
    src = el["src"]
    # Bildpfade sind relativ zum Projekt-Root (REF_DIR/media), nicht zu pages/.
    if os.path.isabs(src):
        path = src
    elif os.path.exists(os.path.join(base_dir, src)):
        path = os.path.join(base_dir, src)
    else:
        path = os.path.join(REF_DIR, src)
    iw, ih = PILImage.open(path).size
    # fit: contain -> Bild vollständig einpassen, zentriert
    scale = min(w / iw, h / ih)
    dw, dh = iw * scale, ih * scale
    dx = x + (w - dw) / 2
    dy = y + (h - dh) / 2
    slide.shapes.add_picture(path, E(dx), E(dy), E(dw), E(dh))


# --------------------------------------------------------------------------- #
#  Folie rendern
# --------------------------------------------------------------------------- #
def render_slide(prs, page_path, colors):
    with open(page_path, "r", encoding="utf-8") as fh:
        page = yaml.safe_load(fh)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # leer

    # Hintergrund: Papier
    bg_color = _resolve_color(page.get("background", {}).get("color", "$paper"))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg_color

    base_dir = os.path.dirname(page_path)
    for el in page.get("elements", []):
        etype = el.get("elementType")
        if etype == "text":
            add_text(slide, el["bounds"], el.get("content", {}))
        elif etype == "shape":
            add_rect(slide, el["bounds"], el.get("fill", {}))
        elif etype == "line":
            add_line(slide, el)
        elif etype == "image":
            add_image(slide, el, base_dir)
        else:
            print("  ! unbekannter Elementtyp:", etype, file=sys.stderr)

    return slide


# --------------------------------------------------------------------------- #
#  Hauptprogramm
# --------------------------------------------------------------------------- #
def main():
    with open(PPTD_PATH, "r", encoding="utf-8") as fh:
        pptd = yaml.safe_load(fh)

    colors = pptd["theme"]["colors"]
    prs = Presentation()
    # 16:9: 13,333" x 7,5"  ==  12.192.000 x 6.858.000 EMU
    prs.slide_width = Emu(960 * 12700)
    prs.slide_height = Emu(540 * 12700)

    cp = pptd["core_properties"] if "core_properties" in pptd else None
    prs.core_properties.title = pptd.get("title", "Die KI-Agenten-Landschaft")

    pages = pptd["pages"]
    print("Rendere %d Folien -> %s" % (len(pages), os.path.relpath(OUT_PPTX, HERE)))
    for i, rel in enumerate(pages, 1):
        page_path = rel if os.path.isabs(rel) else os.path.join(REF_DIR, rel)
        render_slide(prs, page_path, colors)
        print("  [%02d] %s" % (i, os.path.basename(page_path)))

    prs.save(OUT_PPTX)
    print("\nFERTIG: %s  (%d Folien)" % (OUT_PPTX, len(pages)))


if __name__ == "__main__":
    main()
